import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io

st.title("Generador de Presentación")
st.write("Ética del Análisis de Datos Agrícolas")

def crear_presentacion():

    prs = Presentation()

    def fondo(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def titulo(slide, texto):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = texto
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255,255,255)
        p.alignment = PP_ALIGN.CENTER

    def bullets(slide, lista):
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(3))
        tf = box.text_frame
        tf.clear()

        for i,texto in enumerate(lista):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = texto
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255,255,255)

    def placeholder_imagen(slide):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.8),
            Inches(2),
            Inches(3.5),
            Inches(3)
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(255,255,255)

        tf = shape.text_frame
        tf.text = "Imagen aquí"

    def boton(slide,text,left):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            Inches(5),
            Inches(2),
            Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255,255,255)

        tf = shape.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

    # PORTADA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo(slide,RGBColor(34,139,34))
    titulo(slide,"Ética del Análisis de Datos Agrícolas")

    sub = slide.shapes.add_textbox(Inches(2),Inches(2),Inches(6),Inches(1))
    tf=sub.text_frame
    tf.text="Digitalización, privacidad y gobernanza de datos"

    name = slide.shapes.add_textbox(Inches(3),Inches(3),Inches(4),Inches(1))
    tf=name.text_frame
    tf.text="Miguel Ángel Garatejo Capera\nUniversidad del Tolima"

    boton(slide,"Ética",Inches(2))
    boton(slide,"Datos",Inches(4))
    boton(slide,"Agricultura",Inches(6))

    contenido = [

        ("Agricultura Digital",
        ["Uso de sensores y drones",
        "Análisis de datos",
        "Optimización de recursos"]),

        ("Problemas Éticos",
        ["Privacidad de agricultores",
        "Propiedad de datos",
        "Uso por empresas",
        "Falta de transparencia"]),

        ("Privacidad de Datos",
        ["Ubicación de parcelas",
        "Producción agrícola",
        "Uso de insumos",
        "Datos personales"]),

        ("Empresas Tecnológicas",
        ["Plataformas digitales",
        "Acumulación de datos",
        "Dependencia tecnológica"]),

        ("Transparencia",
        ["Qué datos se recolectan",
        "Cómo se usan",
        "Quién accede"]),

        ("Consentimiento Informado",
        ["Comprender el uso",
        "Autorizar datos",
        "Retirar consentimiento"]),

        ("Gobernanza de Datos",
        ["Cooperativas de datos",
        "Data Trusts",
        "Commons digitales"]),

        ("Impactos Sociales",
        ["Mayor productividad",
        "Optimización recursos",
        "Riesgo exclusión"]),

        ("Inclusión",
        ["Mujeres rurales",
        "Jóvenes",
        "Comunidades indígenas"]),

        ("Seguridad de Datos",
        ["Encriptación",
        "Anonimización",
        "Auditorías"]),

        ("Conclusión",
        ["Privacidad",
        "Transparencia",
        "Inclusión",
        "Regulación justa"])
    ]

    for titulo_slide,lista in contenido:

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fondo(slide,RGBColor(60,120,60))

        titulo(slide,titulo_slide)
        bullets(slide,lista)
        placeholder_imagen(slide)

        boton(slide,"Datos",Inches(1))
        boton(slide,"Ética",Inches(3.5))
        boton(slide,"Tecnología",Inches(6))

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer


if st.button("Generar presentación"):

    archivo = crear_presentacion()

    st.download_button(
        label="Descargar PowerPoint",
        data=archivo,
        file_name="exposicion_etica_datos_agricolas.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
